from wand.image import Image
import glob
from pptx import Presentation
from pptx.util import Inches


def ppt_creation():
    j = 0
    prs = Presentation()
    for image_path_name in Output_image_name:
        j = j+1
        image_path= f"image/{image_path_name}.jpg"
        blank_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(blank_slide_layout)

        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = f'Sample title {j}'

        tf = body_shape.text_frame
        tf.text = f'Sample subtitle {j}'

        bottom = Inches(2.5)
        left = Inches(1)
        height = Inches(4)
        image_p = f"{image_path}"
        slide.shapes.add_picture(image_p, left, bottom, height=height)
    prs.save("image/demo.pptx")


def logo_placement(image_val, i):

    logo = Image(filename="image/nike_black.png")
    image = Image(image_val)
    logo_width = int(image.width/2)
    logo_height = int(image.height/6)
    logo.resize(logo_width,logo_height)
    image.composite_channel("all_channels", logo, "dissolve", 70, 30)
    file_save_name = f"outimage{i}"
    image.save(filename=f"image/{file_save_name}.jpg")
    Output_image_name.append(file_save_name)

    ppt_creation()


def get_image():

    i = 0
    for image_name in glob.glob("image/*.jpg"):
        image_val = Image(filename=image_name)
        i = i+1
        logo_placement(image_val, i)


Output_image_name = []
if __name__ == "__main__":
    get_image()




